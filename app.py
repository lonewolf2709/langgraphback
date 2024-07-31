import streamlit as st
from streamlit_extras.add_vertical_space import add_vertical_space
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
import google.generativeai as genai
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import pickle
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
import os
from  dotenv import load_dotenv
from langchain_core.prompts import PromptTemplate
from langchain.chains import LLMChain
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
from agent import agent
from langchain.agents import initialize_agent, Tool
from langchain.chains import RetrievalQA, ConversationalRetrievalChain
from langchain.callbacks.tracers import ConsoleCallbackHandler
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
with st.sidebar:
    st.title('💬 LLM App')
    st.markdown('''
    ## About
    This app is an LLM-powered chatbot built using:
    - [Streamlit](https://streamlit.io/)
    - [LangChain](https://python.langchain.com/)
    - [OpenAI](https://platform.openai.com/docs/models) LLM model
 
    ''')
    add_vertical_space(5)
@st.cache_data(show_spinner="Fetching data...")
def get_summary(docs,language):
    model = ChatGoogleGenerativeAI(model="gemini-pro",temperature=0.5)
    generic_template='''Write a summary of the following text: Text : `{speech}` Translate the precise summary to {language}.'''
    prompt=PromptTemplate(input_variables=['speech','language'],template=generic_template)
    llm_chain=LLMChain(llm=model,prompt=prompt)
    summary=llm_chain.run({'speech':docs,'language':language})
    print(summary)
    return summary
def query_finder(chunks,embeddings,query):
    VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
    docs = VectorStore.similarity_search(query=query, k=18)
    prompt_template = """
                Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
                provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
                Context:\n {context}?\n
                Question: \n{question}\n
    
                Answer:
                """
    model = ChatGoogleGenerativeAI(model="gemini-pro",temperature=0.3)
    prompt = PromptTemplate(template = prompt_template, input_variables = ["context", "question"])
    chain = LLMChain(llm=model,prompt=prompt)
    response = chain({"context":docs, "question": query}, return_only_outputs=True)
    print(response)
    return response["text"]
def read_pdf(file_path):
    text = ""
    reader = PyPDF2.PdfReader(file_path)
    for page in reader.pages:
        text += page.extract_text()
    return text

def read_docx(file_path):
    doc = Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def read_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def read_txt(file_path):
    return file_path.read().decode('utf-8')
def read_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string(index=False)
def read_file(file_path):
    file_type = file_path.name.split('.')[-1].lower()
    print(file_type)
    if file_type == 'pdf':
        return read_pdf(file_path)
    elif file_type == 'docx':
        return read_docx(file_path)
    elif file_type == 'pptx':
        return read_pptx(file_path)
    elif file_type == 'txt':
        return read_txt(file_path)
    elif file_type in ['xls', 'xlsx']:
        return read_excel(file_path)
    else:
        raise ValueError("Unsupported file type")
def main():
    st.header("Chat with PDF 💬")
    file = st.file_uploader("Upload your PDF")
    print(file)
    if file is not None:
        text=read_file(file);
        text_splitter=RecursiveCharacterTextSplitter(
            chunk_size=1500,
            chunk_overlap=350,
            length_function=len
        )
        chunks=text_splitter.split_text(text=text)
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        chunks=text_splitter.split_text(text=text)
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
        VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
        model = ChatGoogleGenerativeAI(model="gemini-pro",temperature=0.5)
        # retriever = VectorStore.as_retriever()
        # qa_chain = load_qa_chain(model, chain_type="stuff")  # 'stuff' is a type of chain for simple QA tasks
        # retrieval_qa_chain = RetrievalQA(combine_documents_chain=qa_chain, retriever=retriever)
        # def query_solver(query):
        #     return retrieval_qa_chain({"query": query})
        # query_solver_tool = Tool(name="query_solver", func=query_solver, description="A tool for answering queries from text.")
        # tools = [Tool(name="query_solver", func=query_solver_tool.run, description="A tool for answering queries from text.")]
        # agent1 = initialize_agent(tools, model,agent_type="zero-shot-react-description", verbose=True)
        query = st.text_input(label="Ask questions about your PDF file:",placeholder="To Summarize the file Write : Summary")
        docs = VectorStore.similarity_search(query=query, k=3)
        if query:
                if(query.lower()=="summary"):
                    args={"input":{"docs":docs,"query":query}}
                    print(type(args),type(args['input']))
                    response= agent.invoke(args,config={"callbacks":[ConsoleCallbackHandler()]})
                    st.write(response['output'])
                else:
                    args={"tool":"Query_Solver","input":{"docs":docs,"query":query}}
                    print(type(args),type(args['input']))
                    response= agent.invoke(args,config={"callbacks":[ConsoleCallbackHandler()]})
                    st.write(response['output'])
if __name__=="__main__":
    main()